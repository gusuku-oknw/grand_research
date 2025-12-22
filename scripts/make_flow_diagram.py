# -*- coding: utf-8 -*-
"""
Flow-style ponchi-e with actors (upload user, search user), servers, shares, and response tiers.
Outputs:
  output/fig_flow.png
  output/fig_flow.pptx
"""
from pathlib import Path
import matplotlib.pyplot as plt
from matplotlib.patches import Rectangle, FancyArrowPatch, Circle
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor

out_dir = Path("output")
out_dir.mkdir(exist_ok=True)

# --- Matplotlib PNG ---
plt.rcParams["font.family"] = ["MS Gothic", "Yu Gothic", "Noto Sans CJK JP", "IPAPGothic", "TakaoPGothic"]
plt.rcParams["axes.unicode_minus"] = False

fig, ax = plt.subplots(figsize=(10, 6))
ax.axis('off')
ax.set_xlim(0, 10)
ax.set_ylim(0, 6)

# Actors positions
actors = {
    "uploader": (1.2, 4.5),
    "searcher": (1.2, 1.5),
    "serverA": (6.0, 4.5),
    "serverB": (6.0, 1.5),
    "recon": (9.0, 3.0)
}

# Draw stick-person icons
for key, (x, y) in actors.items():
    if key in ["serverA", "serverB", "recon"]:
        continue
    head = Circle((x, y+0.4), 0.25, fc='#f2f2f2', ec='black', lw=1.3)
    ax.add_patch(head)
    ax.plot([x, x], [y+0.15, y-0.55], color='black', lw=1.3)
    ax.plot([x, x-0.35], [y-0.2, y-0.7], color='black', lw=1.3)
    ax.plot([x, x+0.35], [y-0.2, y-0.7], color='black', lw=1.3)
    ax.plot([x, x-0.35], [y-0.55, y-1.0], color='black', lw=1.3)
    ax.plot([x, x+0.35], [y-0.55, y-1.0], color='black', lw=1.3)

# Server boxes
server_boxes = {
    "serverA": "シェアサーバA",
    "serverB": "シェアサーバB",
    "recon": "復元ノード (k2到達)"
}
for key, label in server_boxes.items():
    x, y = actors[key]
    rect = Rectangle((x-0.9, y-0.7), 1.8, 1.4, fc='#e8f0ff', ec='black', lw=1.3)
    ax.add_patch(rect)
    ax.text(x, y, label, ha='center', va='center', fontsize=11)

# Shares thumbnails
thumb_colors = {'S1': '#ffd54f', 'S2': '#4fc3f7'}
for i, (name, color) in enumerate(thumb_colors.items()):
    x = 3.0 + i*0.5
    y = 4.7 - i*0.25
    ax.add_patch(Rectangle((x-0.25, y-0.18), 0.5, 0.36, fc=color, ec='black', lw=1))
    ax.text(x, y, f"シェア{name[-1]}", ha='center', va='center', fontsize=9)

# Dummy/original icons
ax.add_patch(Rectangle((7.8, 2.8), 0.7, 0.7, fc='#cfd8dc', ec='black', lw=1))
ax.text(8.15, 3.15, 'ダミー', ha='center', va='center', fontsize=9)
ax.add_patch(Rectangle((8.6, 2.8), 0.7, 0.7, fc='#aed581', ec='black', lw=1))
ax.text(8.95, 3.15, '原画像', ha='center', va='center', fontsize=9)

# Arrows upload
ax.add_patch(FancyArrowPatch( (1.9,4.1), (3.2,4.7), arrowstyle='->', mutation_scale=15, lw=1.3))
ax.text(2.45,4.55,'分割(k1,k2)', fontsize=10)
ax.add_patch(FancyArrowPatch( (3.5,4.7), (5.1,4.5), arrowstyle='->', mutation_scale=15, lw=1.3))
ax.text(4.3,4.9,'シェア1', fontsize=10)
ax.add_patch(FancyArrowPatch( (3.2,4.45), (5.1,1.7), arrowstyle='->', mutation_scale=15, lw=1.3))
ax.text(4.2,3.2,'シェア2', fontsize=10)

# Search arrows
ax.add_patch(FancyArrowPatch( (1.8,1.2), (5.1,1.2), arrowstyle='->', mutation_scale=15, lw=1.3))
ax.text(3.4,1.4,'クエリ(pHash)', fontsize=10)
ax.add_patch(FancyArrowPatch( (6.9,1.2), (8.1,3.0), arrowstyle='->', mutation_scale=15, lw=1.3))
ax.text(7.5,2.0,'k1 未満→ノイズ', fontsize=9)
ax.add_patch(FancyArrowPatch( (6.9,4.5), (8.1,3.0), arrowstyle='->', mutation_scale=15, lw=1.3))
ax.text(7.55,4.0,'k1 達成→ダミー', fontsize=9)
ax.add_patch(FancyArrowPatch( (8.95,3.5), (9.6,3.5), arrowstyle='->', mutation_scale=15, lw=1.3))
ax.text(9.25,3.75,'k2 達成', fontsize=9)

ax.text(1.2,5.4,'アップロード者', ha='center', fontsize=11)
ax.text(1.2,2.1,'検索者', ha='center', fontsize=11)
ax.text(6.0,5.5,'シェア保持サーバ', ha='center', fontsize=11)
ax.text(6.0,2.3,'シェア保持サーバ', ha='center', fontsize=11)
ax.text(9.0,4.1,'応答', ha='center', fontsize=11)

png_path = out_dir / 'fig_flow.png'
fig.savefig(png_path, dpi=200, bbox_inches='tight')
plt.close(fig)

# --- PPTX ---
prs = Presentation()
blank = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank)
slide_width = prs.slide_width
slide_height = prs.slide_height

# helper for person icon
person_width = Inches(0.8)
person_height = Inches(1.6)

def add_person(x_in, y_in, label):
    top = Inches(y_in)
    left = Inches(x_in)
    head = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, top, Inches(0.8), Inches(0.8))
    head.fill.solid(); head.fill.fore_color.rgb = RGBColor(242,242,242); head.line.color.rgb = RGBColor(0,0,0)
    body = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left + Inches(0.25), top+Inches(0.75), Inches(0.3), Inches(0.7))
    body.fill.solid(); body.fill.fore_color.rgb = RGBColor(242,242,242); body.line.color.rgb = RGBColor(0,0,0)
    tb = slide.shapes.add_textbox(left- Inches(0.2), top + Inches(1.6), Inches(1.2), Inches(0.4))
    p = tb.text_frame.paragraphs[0]; p.text = label; p.font.size = Pt(14); p.alignment = 1

add_person(0.6, 1.0, "アップロード者")
add_person(0.6, 4.0, "検索者")

# server boxes
server_specs = [
    (3.8, 0.9, "サーバA (シェア1)"),
    (3.8, 3.9, "サーバB (シェア2)")
]
for x, y, text in server_specs:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.0), Inches(1.2))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(232, 240, 255); shape.line.color.rgb = RGBColor(0,0,0); shape.line.width = Pt(1.5)
    p = shape.text_frame.paragraphs[0]; p.text = text; p.font.size = Pt(14); p.alignment = 1

# recon box
recon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(2.5), Inches(2.2), Inches(1.2))
recon.fill.solid(); recon.fill.fore_color.rgb = RGBColor(224, 242, 232); recon.line.color.rgb = RGBColor(0,0,0); recon.line.width = Pt(1.5)
recon.text_frame.paragraphs[0].text = "応答/復元"
recon.text_frame.paragraphs[0].font.size = Pt(14)
recon.text_frame.paragraphs[0].alignment = 1

# small icons dummy / original
dummy_shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(7.1), Inches(2.6), Inches(0.8), Inches(0.6))
dummy_shape.fill.solid(); dummy_shape.fill.fore_color.rgb = RGBColor(207, 216, 220); dummy_shape.line.color.rgb = RGBColor(0,0,0)
dummy_shape.text_frame.paragraphs[0].text = "ダミー"
dummy_shape.text_frame.paragraphs[0].font.size = Pt(12)
dummy_shape.text_frame.paragraphs[0].alignment = 1
orig_shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(8.0), Inches(2.6), Inches(0.8), Inches(0.6))
orig_shape.fill.solid(); orig_shape.fill.fore_color.rgb = RGBColor(174, 213, 129); orig_shape.line.color.rgb = RGBColor(0,0,0)
orig_shape.text_frame.paragraphs[0].text = "原画像"
orig_shape.text_frame.paragraphs[0].font.size = Pt(12)
orig_shape.text_frame.paragraphs[0].alignment = 1

# connectors
conn1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.4), Inches(1.6), Inches(3.8), Inches(1.4)); conn1.line.end_arrowhead=True; conn1.line.width=Pt(2)
conn1.text = "分割→シェア1"
conn2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.4), Inches(1.4), Inches(3.8), Inches(4.4)); conn2.line.end_arrowhead=True; conn2.line.width=Pt(2); conn2.text="分割→シェア2"
conn3 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.4), Inches(4.6), Inches(3.8), Inches(4.6)); conn3.line.end_arrowhead=True; conn3.line.width=Pt(2); conn3.text="クエリ(pHash)"
conn4 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.8), Inches(1.4), Inches(7.0), Inches(3.1)); conn4.line.end_arrowhead=True; conn4.line.width=Pt(2); conn4.text="k1未満→ノイズ/ダミー"
conn5 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.8), Inches(4.5), Inches(7.0), Inches(3.4)); conn5.line.end_arrowhead=True; conn5.line.width=Pt(2); conn5.text="k1達成→ダミー"
conn6 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.8), Inches(3.1), Inches(9.3), Inches(3.1)); conn6.line.end_arrowhead=True; conn6.line.width=Pt(2); conn6.text="k2達成→原画像"

pptx_path = out_dir / 'fig_flow.pptx'
prs.save(pptx_path)

print('wrote', png_path)
print('wrote', pptx_path)
