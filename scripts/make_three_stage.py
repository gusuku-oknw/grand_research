# -*- coding: utf-8 -*-
"""
Generate three-stage ponchi-e as both PNG and PPTX with editable shapes.
Outputs:
  output/fig_three_stage.png
  output/fig_three_stage.pptx
"""
import math
from pathlib import Path
import matplotlib.pyplot as plt
from matplotlib.patches import Rectangle, FancyArrowPatch
plt.rcParams["font.family"] = ["MS Gothic", "Yu Gothic", "Noto Sans CJK JP", "IPAPGothic", "TakaoPGothic"]
plt.rcParams["font.sans-serif"] = ["MS Gothic", "Yu Gothic", "Noto Sans CJK JP", "IPAPGothic", "TakaoPGothic"]
plt.rcParams["axes.unicode_minus"] = False
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor

out_dir = Path("output")
out_dir.mkdir(exist_ok=True)

# Layout parameters (inches for pptx; also used as axis units for matplotlib)
slide_w_in = 10
slide_h_in = 7.5
box_w = 2.7
box_h = 1.2
center_x = slide_w_in / 2
center_y = 5.8
below_y = 3.3
x_offset = 3.1

positions = {
    "shares": (center_x, center_y),
    "noise": (center_x - x_offset, below_y),
    "dummy": (center_x, below_y),
    "orig": (center_x + x_offset, below_y),
}
texts = {
    "shares": "シェア r 枚\n(受領枚数)",
    "noise": "r < k1\nノイズ画像のみ",
    "dummy": "k1 ≤ r < k2\npHash 符号一致\n視覚はノイズ (ダミー)",
    "orig": "r ≥ k2\n原画像復元",
}
arrows = [
    ("shares", "noise", "閾値 k1 未満"),
    ("shares", "dummy", "閾値 k1 達成"),
    ("shares", "orig", "閾値 k2 達成"),
]
legend = "低周波 DCT 符号のみ一致 → pHash 整合ダミー生成／高周波はノイズで秘匿"

# ---------- PNG via matplotlib ----------
fig, ax = plt.subplots(figsize=(slide_w_in, slide_h_in))
ax.set_xlim(0, slide_w_in)
ax.set_ylim(0, slide_h_in)
ax.axis("off")

for key, (cx, cy) in positions.items():
    x = cx - box_w / 2
    y = cy - box_h / 2
    rect = Rectangle((x, y), box_w, box_h, linewidth=1.5, edgecolor="black", facecolor="#f2f2f2")
    ax.add_patch(rect)
    ax.text(cx, cy, texts[key], ha="center", va="center", fontsize=12)

for src, dst, label in arrows:
    sx, sy = positions[src]
    dx, dy = positions[dst]
    arrow = FancyArrowPatch((sx, sy - box_h/2), (dx, dy + box_h/2),
                            arrowstyle="->", mutation_scale=15, linewidth=1.5)
    ax.add_patch(arrow)
    midx = (sx + dx) / 2
    midy = (sy + dy) / 2 + 0.3
    ax.text(midx, midy, label, ha="center", va="center", fontsize=11)

ax.text(center_x, below_y - 1.2, legend, ha="center", va="center", fontsize=11)

png_path = out_dir / "fig_three_stage.png"
fig.savefig(png_path, dpi=200, bbox_inches="tight")
plt.close(fig)

# ---------- PPTX with editable shapes ----------
prs = Presentation()
prs.slide_width = Inches(slide_w_in)
prs.slide_height = Inches(slide_h_in)
blank = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank)

for key, (cx, cy) in positions.items():
    left = Inches(cx - box_w / 2)
    top = Inches(slide_h_in - cy - box_h / 2)  # pptx origin top-left
    width = Inches(box_w)
    height = Inches(box_h)
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(242, 242, 242)
    shape.line.color.rgb = RGBColor(0, 0, 0)
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.text = texts[key]
    for p in tf.paragraphs:
        p.font.size = Pt(18)
        p.alignment = 1  # center

from pptx.enum.shapes import MSO_CONNECTOR
for src, dst, label in arrows:
    sx, sy = positions[src]
    dx, dy = positions[dst]
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Inches(sx), Inches(slide_h_in - sy + box_h/2),
                                      Inches(dx), Inches(slide_h_in - dy - box_h/2))
    conn.line.end_arrowhead = True
    conn.line.width = Pt(2)
    midx = (sx + dx) / 2
    midy = (sy + dy) / 2
    tb = slide.shapes.add_textbox(Inches(midx - 1.2), Inches(slide_h_in - midy - 0.2), Inches(2.4), Inches(0.5))
    para = tb.text_frame.paragraphs[0]
    para.text = label
    para.font.size = Pt(16)
    para.alignment = 1

leg_tb = slide.shapes.add_textbox(Inches(center_x - 4.2), Inches(slide_h_in - (below_y - 1.2) - 0.25), Inches(8.4), Inches(0.6))
leg_p = leg_tb.text_frame.paragraphs[0]
leg_p.text = legend
leg_p.font.size = Pt(16)
leg_p.alignment = 1

pptx_path = out_dir / "fig_three_stage.pptx"
prs.save(pptx_path)

print("wrote", png_path)
print("wrote", pptx_path)
