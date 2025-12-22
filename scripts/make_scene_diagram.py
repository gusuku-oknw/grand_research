# -*- coding: utf-8 -*-
"""
Scene-style ponchi-e with actors + photo + shares + servers + responses.
Outputs:
  output/fig_scene.png
  output/fig_scene.pptx
"""
from pathlib import Path
import matplotlib.pyplot as plt
from matplotlib.patches import Rectangle, Circle, FancyArrowPatch
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor

out_dir = Path("output")
out_dir.mkdir(exist_ok=True)

plt.rcParams["font.family"] = ["MS Gothic", "Yu Gothic", "Noto Sans CJK JP", "IPAPGothic", "TakaoPGothic"]
plt.rcParams["axes.unicode_minus"] = False

fig, ax = plt.subplots(figsize=(10.5, 5.8))
ax.axis('off')
ax.set_xlim(0, 10)
ax.set_ylim(0, 5.5)

# simple round avatars (視認性重視)
def avatar(x, y, label, color):
    circ = Circle((x, y), 0.32, fc=color, ec='black', lw=1.2)
    ax.add_patch(circ)
    ax.text(x, y, label, ha='center', va='center', fontsize=11, color='white', weight='bold')

# avatar placements
avatar(1.0, 4.5, "U", "#42a5f5")   # uploader
avatar(1.0, 1.6, "Q", "#ffb74d")   # querier/searcher
ax.text(1.0, 4.05, "撮影者/配布者", ha='center', va='top', fontsize=10)
ax.text(1.0, 1.1, "検索者", ha='center', va='top', fontsize=10)

# photo icon
ax.add_patch(Rectangle((1.55,4.3), 0.95,0.65, fc='#fff3cd', ec='black', lw=1.2))
ax.text(2.02,4.6,'画像', ha='center', va='center', fontsize=10)

# share squares
share_colors=['#ffd54f','#4fc3f7','#ce93d8','#80cbc4']
for i,c in enumerate(share_colors):
    ax.add_patch(Rectangle((2.85+i*0.5,4.15-0.1*i),0.45,0.32, fc=c, ec='black', lw=1))
    ax.text(3.08+i*0.5,4.31-0.1*i,f'S{i+1}', ha='center', va='center', fontsize=9)

# servers
server_specs=[(5.3,4.1,'シェア保持サーバA(S1)'),
              (5.3,2.7,'シェア保持サーバB(S2)'),
              (5.3,1.3,'シェア保持サーバC(S3,S4)')]
for x,y,name in server_specs:
    ax.add_patch(Rectangle((x-0.9,y-0.5),1.8,1.1, fc='#e8f0ff', ec='black', lw=1.2))
    ax.text(x, y+0.05, name, ha='center', va='center', fontsize=10)

# response panel
ax.add_patch(Rectangle((8.1,2.1),1.55,1.55, fc='#e0f2f1', ec='black', lw=1.2))
ax.text(8.88,3.45,'応答', ha='center', fontsize=11)
ax.add_patch(Rectangle((8.25,2.8),0.55,0.42, fc='#cfd8dc', ec='black', lw=1))
ax.text(8.52,3.01,'ダミー', ha='center', fontsize=8)
ax.add_patch(Rectangle((9.05,2.8),0.55,0.42, fc='#aed581', ec='black', lw=1))
ax.text(9.32,3.01,'原画像', ha='center', fontsize=8)
ax.text(8.88,2.45,'k1未満→ノイズ\nk1到達→ダミー\nk2到達→原画像', ha='center', va='top', fontsize=9)

# arrows upload
ax.add_patch(FancyArrowPatch((1.5,4.5),(2.65,4.5), arrowstyle='->', mutation_scale=14, lw=1.2))
ax.text(2.05,4.72,'pHash + 分割', fontsize=9, ha='center')
ax.add_patch(FancyArrowPatch((3.3,4.35),(4.4,4.35), arrowstyle='->', mutation_scale=14, lw=1.2))
ax.text(3.85,4.55,'S1', fontsize=9)
ax.add_patch(FancyArrowPatch((3.8,4.15),(4.4,2.95), arrowstyle='->', mutation_scale=14, lw=1.2))
ax.text(4.25,3.6,'S2', fontsize=9)
ax.add_patch(FancyArrowPatch((4.3,3.95),(4.4,1.75), arrowstyle='->', mutation_scale=14, lw=1.2))
ax.text(4.35,2.85,'S3/S4', fontsize=9)

# search arrows
ax.add_patch(FancyArrowPatch((1.45,1.9),(4.4,1.9), arrowstyle='->', mutation_scale=14, lw=1.2))
ax.text(2.95,2.1,'クエリ(pHash)', fontsize=9, ha='center')
ax.add_patch(FancyArrowPatch((6.05,1.9),(7.95,2.9), arrowstyle='->', mutation_scale=14, lw=1.2))
ax.text(6.95,2.45,'k1～k2 集約', fontsize=9, ha='center')
ax.add_patch(FancyArrowPatch((8.0,2.9),(9.85,2.9), arrowstyle='->', mutation_scale=14, lw=1.2))
ax.text(8.95,3.15,'閲覧結果', fontsize=9, ha='center')

png_path = out_dir/'fig_scene.png'
fig.savefig(png_path, dpi=220, bbox_inches='tight')
plt.close(fig)

# --- PPTX version ---
prs = Presentation()
blank = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank)
slide_width = prs.slide_width
slide_height = prs.slide_height

# helpers

def add_person_ppt(x_in, y_in, label):
    head = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x_in), Inches(y_in), Inches(0.7), Inches(0.7))
    head.fill.solid(); head.fill.fore_color.rgb = RGBColor(242,242,242); head.line.color.rgb = RGBColor(0,0,0)
    body = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x_in+0.25), Inches(y_in+0.7), Inches(0.2), Inches(0.7))
    body.fill.solid(); body.fill.fore_color.rgb = RGBColor(242,242,242); body.line.color.rgb = RGBColor(0,0,0)
    tb = slide.shapes.add_textbox(Inches(x_in-0.2), Inches(y_in+1.5), Inches(1.2), Inches(0.4))
    p = tb.text_frame.paragraphs[0]; p.text = label; p.font.size = Pt(14); p.alignment = 1

def add_server(x_in, y_in, text):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x_in), Inches(y_in), Inches(2.0), Inches(1.0))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(232,240,255); shape.line.color.rgb = RGBColor(0,0,0); shape.line.width = Pt(1.5)
    p = shape.text_frame.paragraphs[0]; p.text = text; p.font.size = Pt(13); p.alignment = 1

add_person_ppt(0.5, 0.7, "撮影者/アップロード")
add_person_ppt(0.5, 3.2, "検索者")
add_server(3.0, 0.7, "サーバA (シェア1)")
add_server(3.0, 3.2, "サーバB (シェア2)")

# response box
resp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(1.8), Inches(2.5), Inches(1.5))
resp.fill.solid(); resp.fill.fore_color.rgb = RGBColor(224,242,232); resp.line.color.rgb = RGBColor(0,0,0); resp.line.width = Pt(1.5)
resp.text_frame.paragraphs[0].text = "応答\nk1未満:ノイズ\nk1到達:ダミー\nk2到達:原画像"
resp.text_frame.paragraphs[0].font.size = Pt(13)
resp.text_frame.paragraphs[0].alignment = 1

# arrows
conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.2), Inches(1.0), Inches(3.0), Inches(1.0)); conn.line.end_arrowhead=True; conn.line.width=Pt(2); conn.text="pHash+分割"
conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.2), Inches(1.0), Inches(3.0), Inches(3.7)); conn.line.end_arrowhead=True; conn.line.width=Pt(2); conn.text="S2"
conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.2), Inches(3.9), Inches(3.0), Inches(3.9)); conn.line.end_arrowhead=True; conn.line.width=Pt(2); conn.text="クエリ(pHash)"
conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.0), Inches(1.0), Inches(6.5), Inches(2.55)); conn.line.end_arrowhead=True; conn.line.width=Pt(2)
conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.0), Inches(4.2), Inches(6.5), Inches(2.85)); conn.line.end_arrowhead=True; conn.line.width=Pt(2)
conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.9), Inches(2.55), Inches(9.5), Inches(2.55)); conn.line.end_arrowhead=True; conn.line.width=Pt(2); conn.text="閲覧"

pptx_path = out_dir/'fig_scene.pptx'
prs.save(pptx_path)

print('wrote', png_path)
print('wrote', pptx_path)
